import math
import os
import random
import secrets

import numpy as np
import win32com.client
from pptx.dml.color import RGBColor
from pptx.text.fonts import FontFiles
from pptx.util import Pt
from tqdm import tqdm

import config
import utils


def generation(markov_text, fonts, backgrounds):
    # create an intermediate presentation, from which we will receive the coordinates of the bboxes, and the final
    prs, slide = utils.make_prs(config.SLIDE_WIDTH, config.SLIDE_HEIGHT)
    final_prs, final_slide = utils.make_prs(config.SLIDE_WIDTH, config.SLIDE_HEIGHT)
    # generate a unique name for the sample
    unic_name = secrets.token_hex(nbytes=16)
    prs_path = os.path.join(config.TEMP_DIR, 'temp.pptx')
    final_prs_path = os.path.join(config.OUTPUT_DIR, 'synthpptx_' + unic_name + '.pptx')
    screenshot_path = os.path.join(config.TEMP_DIR, 'temp.jpg')
    txt_path = os.path.join(config.OUTPUT_DIR, 'synthpptx_' + unic_name + '.txt')

    type_of_background = random.random()
    # randomly insert a color or picture background, or leave the default white
    if type_of_background < config.COLOUR_BACKGROUND_PROB:
        background_fill = final_slide.background.fill
        background_fill.solid()
        background_fill.fore_color.rgb = RGBColor(*utils.get_random_color())
    elif type_of_background < config.COLOUR_BACKGROUND_PROB + config.PICTURE_BACKGROUND_PROB:
        background_image = random.sample(backgrounds, 1)[0]
        pic = final_slide.shapes.add_picture(background_image, left=Pt(0), top=Pt(0),
                                             height=Pt(config.SLIDE_HEIGHT), width=Pt(config.SLIDE_WIDTH))

    # randomly initialize line spacing for all text on a slide
    random_margin_after_lines = random.randint(0, config.MAX_MARGIN_AFTER_LINES)

    # randomly initialize parameters for all content_titles on a slide
    content_titles_color = RGBColor(*utils.get_simple_random_color())
    content_titles_font = random.sample(fonts, 1)[0][0]
    content_titles_font_size = random.randint(config.MIN_STRING_SIZE, config.MAX_STRING_SIZE)
    bold_content_titles = True if random.random() < config.CONTENT_TITLE_BOLD_PROB else False
    upper_content_titles = True if random.random() < config.CONTENT_TITLE_UPPER_PROB else False

    # randomly initialize the parameters for all bullets on the slide
    bullets_color = RGBColor(*utils.get_random_color())
    bullets_font = random.sample(fonts, 1)[0][0]
    bullets_font_size = min(content_titles_font_size, random.randint(config.MIN_STRING_SIZE, config.MAX_STRING_SIZE))
    bold_bullets = True if random.random() < config.BULLET_BOLD_PROB and bold_content_titles else False
    upper_bullets = True if random.random() < config.BULLET_UPPER_PROB and upper_content_titles else False

    # randomly initialize the number of content_blocks on the slide and the vertical distance between them
    block_count = random.randint(config.MIN_BLOCK_COUNT, config.MAX_BLOCK_COUNT)
    blocks_vertical_gap = random.randint(0, content_titles_font_size + random_margin_after_lines)

    # randomly initialize distances between content_title and bullets
    vertical_class_gap = random.randint(-random_margin_after_lines, random_margin_after_lines)
    horizontal_class_gap = config.HORIZONTAL_CLASS_GAP_COEFF * random.randint(-random_margin_after_lines,
                                                                              random_margin_after_lines)
    # randomly initialize the interval between adjacent bullets for all bullets on the slide
    bullets_gap = random.randint(2 * random_margin_after_lines,
                                 max(2 * random_margin_after_lines, math.ceil(bullets_font_size / 2)))

    slide_array = []
    for b, block in enumerate(range(block_count)):
        # for the first block, select a suitable random point in the upper half of the slide
        if b == 0:
            content_title_left = random.randint(0, config.SLIDE_WIDTH - 1050)
            content_title_top = random.randint(150, config.SLIDE_HEIGHT - 600)
        else:
            content_title_top = bullet_bbox[3] + blocks_vertical_gap

        # samples a list of a random number of lines with random text
        content_titles_lines = random.randint(1, config.MAX_CONTENT_TITLE_LINES)
        content_titles_text = utils.gen_text(markov_text, content_titles_lines, config.MIN_WORDS,
                                             config.MAX_WORDS, upper_content_titles)

        # randomly add a colon at the end of the content_title
        if random.random() < config.CONTENT_TITLE_COLON_PROB:
            content_titles_text[-1] += ':'

        # sequentially insert content_title lines on the intermediate and final slide
        for t, text_line in enumerate(content_titles_text):
            utils.add_textbox(slide, content_title_left, content_title_top,
                              config.TEXT_SHAPE_WIDTH, config.TEXT_SHAPE_HEIGHT, content_titles_font,
                              content_titles_font_size, content_titles_color, text_line, bold_content_titles)
            utils.add_textbox(final_slide, content_title_left, content_title_top,
                              config.TEXT_SHAPE_WIDTH, config.TEXT_SHAPE_HEIGHT, content_titles_font,
                              content_titles_font_size, content_titles_color, text_line, bold_content_titles)
            content_title_top += content_titles_font_size + random_margin_after_lines

        # save the intermediate presentation and its screenshot
        prs.save(prs_path)
        pres = PptApp.Presentations.Open(prs_path)
        pres.Slides[0].Export(screenshot_path, "JPG")
        pres.close()

        # calculate the coordinates of the content_title bbox, add them to the general array of coordinates
        content_title_bbox, screenshot = utils.get_bbox(screenshot_path, config.SLIDE_WIDTH, config.SLIDE_HEIGHT)
        content_title_bbox.append(1)
        slide_array.append(content_title_bbox)

        # make alltext on the intermediate slide white
        # so that it does not interfere with further calculation of coordinates
        for sh in slide.shapes:
            sh.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

        # create a content_block with coordinates of its content_title (later we will expand the block)
        block_bbox = content_title_bbox[:-1]
        block_bbox.append(2)

        # randomly initialize the number of bullets in the current block
        bullets_count = random.randint(1, config.MAX_BULLET_COUNT)
        for j, bullet in enumerate(range(bullets_count)):
            # move the first bullet in the block from the content_title to the vertical_class_gap distance
            if j == 0:
                bullet_left = content_title_bbox[0] + horizontal_class_gap
                bullet_top = content_title_bbox[3] + vertical_class_gap
            else:
                bullet_top = bullet_bbox[3] + bullets_gap

            # samples a list of a random number of lines with random text
            bullets_lines = random.randint(1, config.MAX_BULLET_LINES)
            bullets_text = utils.gen_text(markov_text, bullets_lines, config.MIN_WORDS, config.MAX_WORDS, upper_bullets)

            # randomly add a red line at the beginning of the bullet
            if random.random() < config.RED_LINE_PROB:
                bullets_text[0] = '    ' + bullets_text[0][:-4]

            # sequentially insert the lines of the bullet on the intermediate and final slide
            for t, text_line in enumerate(bullets_text):
                utils.add_textbox(slide, bullet_left, bullet_top,
                                  config.TEXT_SHAPE_WIDTH, config.TEXT_SHAPE_HEIGHT, bullets_font,
                                  bullets_font_size, bullets_color, text_line, bold_bullets)
            utils.add_textbox(final_slide, bullet_left, bullet_top,
                              config.TEXT_SHAPE_WIDTH, config.TEXT_SHAPE_HEIGHT, bullets_font,
                              bullets_font_size, bullets_color, text_line, bold_bullets)
            bullet_top += bullets_font_size + random_margin_after_lines

            # save the intermediate presentation and its screenshot
            prs.save(prs_path)
            pres = PptApp.Presentations.Open(prs_path)
            pres.Slides[0].Export(screenshot_path, "JPG")
            pres.close()

            # calculate the coordinates of the bbox bullet, add them to the general array of coordinates
            bullet_bbox, screenshot = utils.get_bbox(screenshot_path, config.SLIDE_WIDTH, config.SLIDE_HEIGHT)
            bullet_bbox.append(3)
            slide_array.append(bullet_bbox)

            # make alltext on the intermediate slide white
            # so that it does not interfere with further calculation of coordinates
            for sh in slide.shapes:
                sh.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

            # update the coordinates of the current content_block
            block_bbox = [min(block_bbox[0], bullet_bbox[0]), min(block_bbox[1], bullet_bbox[1]),
                          max(block_bbox[2], bullet_bbox[2]), max(block_bbox[3], bullet_bbox[3]), 2]

        # add the current content block to the general array
        slide_array.append(block_bbox)

        # make alltext on the intermediate slide white
        # so that it does not interfere with further calculation of coordinates
        for sh in slide.shapes:
            sh.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # remove possible bboxes outside the slide -
    # they have coordinates [SLIDE_WIDTH*SLIDE_HEIGHT, SLIDE_WIDTH*SLIDE_HEIGHT, 0, 0])
    slide_array = [row for row in slide_array if config.SLIDE_WIDTH*config.SLIDE_HEIGHT not in row]

    # calculate and add a block for content_blocks to a common array
    slide_array = np.array(slide_array)
    general_bbox = [np.min(slide_array[:, 0]), np.min(slide_array[:, 1]),
                    np.max(slide_array[:, 2]), np.max(slide_array[:, 3]), 2]
    slide_array = np.append(slide_array, [general_bbox], axis=0)

    # save the markup and the final slide
    utils.save_txt(txt_path, slide_array, config.SLIDE_WIDTH, config.SLIDE_HEIGHT)
    final_prs.save(final_prs_path)

    # save a screenshot of the final presentation
    pres = PptApp.Presentations.Open(final_prs_path)
    pres.Slides[0].Export(final_prs_path[:-4] + 'jpg', "JPG")
    pres.close()

    # draw bbox visualization over the last screenshot (only for check final bboxes quality)
    visualization = utils.draw_visualization(screenshot, slide_array)

    # break


if __name__ == "__main__":
    markov_text = utils.generate_random_text(config.TEXT_DONOR, config.TEXT_SAMPLES)

    # load client for screenshotting slides
    PptApp = win32com.client.Dispatch("Powerpoint.Application")
    PptApp.Visible = True

    FONTS = [f for f in list(FontFiles._installed_fonts()) if f[0] not in config.BAD_FONTS]  # loading installed fonts

    BACKGROUNDS_LIST = utils.get_pictures(config.BACKGROUNDS_DIR)  # creating a list of backgrounds to sample from it
    for _ in tqdm(range(config.PPTX_COUNT)):
        generation(markov_text, FONTS, BACKGROUNDS_LIST)
