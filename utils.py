import glob
import os
import random

import cv2
import markovgen
import numpy as np
import pythoncom
import win32com.client
from pptx import Presentation
from pptx.util import Pt
from tqdm import tqdm

import config


def generate_random_text(textfile, n):
    """Generate long text samples, remove quote characters from them."""
    text_donor = open(textfile)
    markov = markovgen.Markov(text_donor)
    markov_text = list(set([markov.generate_markov_text(max_size=1000, backward=True).replace('"', '').replace("'", '')
                            for _ in tqdm(range(n))]))
    return markov_text


def get_pictures(root_dir):
    """Paths to all images in the folder."""
    paths_to_images = glob.glob(root_dir + '/**/*.jpg', recursive=True)
    return paths_to_images


def image_resize(image, width=None, height=None, inter=cv2.INTER_AREA):
    """Resize image to given height/width."""
    dim = None
    (h, w) = image.shape[:2]

    # if both the width and height are None, then return the
    # original image
    if width is None and height is None:
        return image

    # check to see if the width is None
    if width is None:
        r = height / float(h)
        dim = (int(w * r), height)

    # otherwise, the height is None
    else:
        r = width / float(w)
        dim = (width, int(h * r))

    # resize the image
    resized = cv2.resize(image, dim, interpolation=inter)

    # return the resized image
    return resized


def get_random_color(max_value=255):
    """Random color generation."""
    r = lambda: random.randint(0, max_value)
    random_color = ('#%02X%02X%02X' % (r(), r(), r()))
    return random_color


def get_simple_random_color():
    """Generating a random dark color from a list of the most popular font colors."""
    r = random.sample([(0, 0, 0), (54, 69, 79), (2, 48, 32), (48, 25, 52), (52, 52, 52),
                       (27, 18, 18), (40, 40, 43), (25, 25, 112), (53, 57, 53)], 1)[0]
    random_color = ('#%02X%02X%02X' % r)
    return random_color


def first_nonzero(arr, axis, threshed=False, mask_val=0, invalid_val=-1):
    """Getting the first non-null element in the array along the selected axis."""

    # depending on the arguments, choose what to consider as the background - black or threshold color
    if not threshed:
        mask = arr != mask_val
    else:
        mask = arr > mask_val
    val = np.where(mask.any(axis=axis), mask.argmax(axis=axis), invalid_val)
    return np.where(mask.any(axis=axis), val, invalid_val)


def last_nonzero(arr, axis, threshed=False, mask_val=0, invalid_val=-1):
    """Getting the last non-null element in the array along the selected axis."""

    # depending on the arguments, choose what to consider as the background - black or threshold color
    if not threshed:
        mask = arr != mask_val
    else:
        mask = arr > mask_val
    val = arr.shape[axis] - np.flip(mask, axis=axis).argmax(axis=axis) - 1
    return np.where(mask.any(axis=axis), val, invalid_val)


def make_prs(slide_width, slide_height):
    """Creating a presentation from a single slide of a given width and height."""
    prs = Presentation()
    prs.slide_width, prs.slide_height = Pt(slide_width), Pt(slide_height)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    return prs, slide


def gen_text(markov_text, lines, min_words, max_words, is_upper):
    """Generation of text for substitution in bbox."""
    text = random.sample(markov_text, lines)
    max_words = random.randint(min_words, max_words)
    # sort the text by length so that the shortest line is always at the bottom
    text = sorted(text, key=lambda x: len(x), reverse=True)
    # just in case, make the text in the first line long
    text[0] = '. '.join(text)
    # we leave only max_words words in each line
    text = [' '.join(t.split()[:max_words]) for t in text]
    # make all lines no longer than the first
    text = [t[:len(text[0])] for t in text]

    # randomly convert all text or the first characters in strings to uppercase
    if is_upper:
        text = [s.upper() for s in text]
    else:
        for j, sentence in enumerate(text):
            # the first line is always capitalized
            if j != 0:
                # convert first characters in random strings to lower case
                if random.random() > 0.5:
                    text[j] = sentence[0].lower() + sentence[1:]
    return text


def add_textbox(slide, left, top, width, height, font, font_size, font_color, text_line, is_bold):
    """Adding a text shape-string with the specified parameters to a slide."""
    textbox = slide.shapes.add_textbox(Pt(left), Pt(top), Pt(width), Pt(height))
    text_frame = textbox.text_frame
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    paragraph.font.name = font
    paragraph.font.size = Pt(font_size)
    paragraph.font.color.rgb = font_color
    paragraph.text = text_line
    paragraph.font.bold = is_bold


def get_bbox(screenshot_path, slide_width, slide_height):
    """Getting bbox text on slide screenshot."""
    image = cv2.imread(screenshot_path)
    # invert the image so the background is black
    image = cv2.bitwise_not(image)
    # resize the image to the size of the slide
    # (because for some reason a screenshot of 1200 * 900 of the slide turns out to be 1600 * 1200 in size)
    image = image_resize(image, slide_width, slide_height)
    # for top_left point, we fill array elements responsible for background pixels with this very large number
    square = slide_width * slide_height
    # find the boundaries of the bbox by the first/last mismatch of values with the values of the argument invalid_val
    left = max(0, np.min(first_nonzero(image, 1, invalid_val=square)))
    top = max(0, np.min(first_nonzero(image, 0, invalid_val=square)))
    right = min(image.shape[1], np.max(last_nonzero(image, 1, invalid_val=-1)) + 1)
    bottom = min(image.shape[0], np.max(last_nonzero(image, 0, invalid_val=-1)) + 1)
    return [left, top, right, bottom], cv2.bitwise_not(image)


def save_txt(txt_path, slide_array, slide_width, slide_height):
    """Save markup in yolo format."""
    with open(txt_path, 'a') as txt:
        for ann in slide_array:
            label = ann[-1]
            x1 = ann[0]
            y1 = ann[1]
            x2 = ann[2]
            y2 = ann[3]
            xmax, xmin = x2, x1
            ymax, ymin = y2, y1
            dw = 1./slide_width
            dh = 1./slide_height
            x = (xmin + xmax)/2.0
            y = (ymin + ymax)/2.0
            w = xmax - xmin
            h = ymax - ymin
            x = x*dw
            w = w*dw
            y = y*dh
            h = h*dh
            bb = (x, y, w, h)
            txt.write(str(label) + " " + " ".join([str(a) for a in bb]) + '\n')


def draw_visualization(background, array):
    """Drawing bboxes on the background."""
    visualization = np.copy(background)
    for bbox in array:
        x1 = bbox[0]
        y1 = bbox[1]
        x2 = bbox[2]
        y2 = bbox[3]
        label = bbox[-1]
        color = config.LABELS_TO_COLORS[config.ALL_CLASSES[label]]
        visualization = cv2.rectangle(visualization, (x1, y1), (x2, y2), color, 2)
        cv2.putText(visualization, config.ALL_CLASSES[label], (x1, y1), cv2.FONT_HERSHEY_SIMPLEX, 0.3, (0, 0, 0), 1)
    return visualization


def ppt2pptx(folder, out_folder):
    """Convert PPT to PPTX via win32com.client."""

    PptApp = win32com.client.Dispatch("Powerpoint.Application")
    PptApp.Visible = True

    for prs_path in tqdm(glob.glob(folder + '/**/*.ppt', recursive=True)):
        try:
            pres = PptApp.Presentations.Open(prs_path)
            pres.SaveAs(os.path.join(out_folder, os.path.basename(prs_path) + 'x'), 24)
            pres.close()
        except pythoncom.com_error:
            print(f'skip {prs_path}')
            continue


def pptx2jpg(folder, out_folder):
    """Convert PPTX to JPG via win32com.client."""

    PptApp = win32com.client.Dispatch("Powerpoint.Application")
    PptApp.Visible = True

    for prs_path in tqdm(glob.glob(folder + '/**/*.pptx', recursive=True)):
        pres = PptApp.Presentations.Open(prs_path)
        screenshot_path = os.path.join(out_folder, prs_path[:-4] + 'jpg')
        pres.Slides[0].Export(screenshot_path, "JPG")
        pres.close()

        image = cv2.imread(screenshot_path)
        image = image_resize(image, 1200, 900)
        cv2.imwrite(screenshot_path, image, [cv2.IMWRITE_JPEG_QUALITY, 100])
